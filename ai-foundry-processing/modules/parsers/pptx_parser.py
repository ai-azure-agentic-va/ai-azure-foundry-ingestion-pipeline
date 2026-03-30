"""PowerPoint parser using python-pptx - no Azure Document Intelligence."""

import io
import logging
from pptx import Presentation
from .base import BaseParser, ParseResult

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]

    def parse(self, file_bytes: bytes) -> ParseResult:
        logger.info(f"[PptxParser] Parsing PPTX ({len(file_bytes)} bytes)")
        prs = Presentation(io.BytesIO(file_bytes))

        slides = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(row_text))

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"[Speaker Notes] {notes}")

            if texts:
                slides.append(
                    {
                        "page_number": slide_num,
                        "text": "\n".join(texts),
                    }
                )

        full_text = "\n\n".join(
            f"Slide {s['page_number']}:\n{s['text']}" for s in slides
        )

        logger.info(
            f"[PptxParser] Extracted {len(slides)} slides, {len(full_text)} chars"
        )
        return ParseResult(
            full_text=full_text,
            pages=[
                {"page_number": s["page_number"], "text": s["text"], "table_text": ""}
                for s in slides
            ],
            page_count=len(slides),
            metadata={"format": "pptx"},
        )
